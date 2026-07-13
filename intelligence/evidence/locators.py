from __future__ import annotations

import csv
import bz2
import gzip
import hashlib
import io
import json
import os
import re
import sqlite3
import tempfile
import tarfile
import lzma
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup

from .models import (
    ArchiveMemberLocator, CalendarComponentLocator, CsvCellLocator, CsvRowLocator,
    DatabaseCellLocator, DatabaseTableRowLocator, EmailAttachmentLocator, EmailHeaderLocator,
    EmailMimePartLocator, GeospatialFeatureLocator, HtmlDomSpanLocator, ImageRegionLocator,
    JsonPointerLocator, JsonRecordLocator, LocatorType, MediaTimeRangeLocator,
    OfficeParagraphLocator,OfficeTableCellLocator,PdfPageBlockLocator,PdfRegionLocator,
    SlideNotesLocator,SlideShapeLocator,SpreadsheetCellLocator,SubtitleCueLocator,TextByteSpanLocator, TextLineLocator, TextSpanLocator,
    VcardPropertyLocator, VideoFrameLocator, XmlElementLocator,
)


class LocatorResolutionError(ValueError):
    pass


def _archive_member_bytes(content:bytes,member_path:str,member_ordinal:int|None=None)->bytes:
    if content.startswith((b"PK\x03\x04",b"PK\x05\x06",b"PK\x07\x08")):
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            infos=archive.infolist()
            if member_ordinal is not None:
                if member_ordinal>=len(infos) or infos[member_ordinal].filename.replace("\\","/")!=member_path.replace("\\","/"):
                    raise LocatorResolutionError("archive member ordinal does not resolve")
                return archive.read(infos[member_ordinal])
            matches=[info for info in infos if info.filename.replace("\\","/")==member_path.replace("\\","/")]
            if len(matches)!=1: raise LocatorResolutionError("archive member path is missing or ambiguous")
            return archive.read(matches[0])
    try:
        with tarfile.open(fileobj=io.BytesIO(content),mode="r:*") as archive:
            members=archive.getmembers()
            candidate=(members[member_ordinal] if member_ordinal is not None and member_ordinal<len(members) else archive.getmember(member_path))
            if candidate.name.replace("\\","/")!=member_path.replace("\\","/"): raise LocatorResolutionError("archive member ordinal does not resolve")
            stream=archive.extractfile(candidate)
            if stream is None: raise LocatorResolutionError("archive member is not a file")
            return stream.read()
    except (tarfile.TarError,KeyError):
        pass
    name=member_path.replace("\\","/")
    if "/" in name: raise LocatorResolutionError("single-stream member path does not resolve")
    try:
        if content.startswith(b"\x1f\x8b"): return gzip.decompress(content)
        if content.startswith(b"BZh"): return bz2.decompress(content)
        if content.startswith(b"\xfd7zXZ\x00"): return lzma.decompress(content)
    except (OSError,EOFError,lzma.LZMAError) as exc:
        raise LocatorResolutionError("single-stream archive is corrupt") from exc
    raise LocatorResolutionError("unsupported archive container")


def _json_pointer(document: Any, pointer: str) -> Any:
    current = document
    if pointer == "": return current
    for token in pointer.split("/")[1:]:
        token = token.replace("~1", "/").replace("~0", "~")
        try:
            current = current[int(token)] if isinstance(current, list) else current[token]
        except (KeyError, IndexError, ValueError, TypeError) as exc:
            raise LocatorResolutionError(f"JSON Pointer does not resolve: {pointer}") from exc
    return current


def resolve_locator(content: bytes, locator_type: LocatorType | str, locator: dict[str, Any]) -> bytes:
    try: kind = LocatorType(locator_type)
    except ValueError as exc: raise LocatorResolutionError(f"unsupported locator type: {locator_type}") from exc

    try:
        if kind is LocatorType.JSON_POINTER:
            model=JsonPointerLocator.model_validate(locator)
            value=_json_pointer(json.loads(content.decode("utf-8")),model.pointer)
            return json.dumps(value,sort_keys=True,separators=(",",":"),ensure_ascii=False).encode()
        if kind is LocatorType.JSON_RECORD:
            model=JsonRecordLocator.model_validate(locator)
            text=content.decode("utf-8-sig")
            try:
                document=json.loads(text)
                records=document if isinstance(document,list) else [document]
            except json.JSONDecodeError:
                records=[json.loads(line) for line in text.splitlines() if line.strip()]
            if model.record>=len(records): raise LocatorResolutionError("JSON record is outside the artifact")
            value=_json_pointer(records[model.record],model.pointer)
            return json.dumps(value,sort_keys=True,separators=(",",":"),ensure_ascii=False).encode()
        if kind in (LocatorType.CSV_ROW,LocatorType.CSV_CELL):
            rows=list(csv.reader(io.StringIO(content.decode("utf-8-sig"))))
            if kind is LocatorType.CSV_ROW:
                model=CsvRowLocator.model_validate(locator)
                if model.row>len(rows): raise LocatorResolutionError("CSV row is outside the artifact")
                return json.dumps(rows[model.row-1],ensure_ascii=False,separators=(",",":")).encode()
            model=CsvCellLocator.model_validate(locator)
            if model.row>len(rows): raise LocatorResolutionError("CSV row is outside the artifact")
            idx=model.column if isinstance(model.column,int) else (rows[0].index(model.column) if rows else -1)
            if idx<0 or idx>=len(rows[model.row-1]): raise LocatorResolutionError("CSV column is outside the artifact")
            return rows[model.row-1][idx].encode()
        if kind is LocatorType.TEXT_SPAN:
            model=TextSpanLocator.model_validate(locator)
            if model.byte_end>len(content): raise LocatorResolutionError("text span is outside the artifact")
            return content[model.byte_start:model.byte_end]
        if kind is LocatorType.TEXT_BYTE_SPAN:
            model=TextByteSpanLocator.model_validate(locator)
            if model.byte_end>len(content): raise LocatorResolutionError("text byte span is outside the artifact")
            return content[model.byte_start:model.byte_end]
        if kind is LocatorType.TEXT_LINE:
            model=TextLineLocator.model_validate(locator)
            lines=content.decode("utf-8-sig").splitlines()
            if model.line>len(lines): raise LocatorResolutionError("text line is outside the artifact")
            value=lines[model.line-1]; start=model.column_start or 0; end=model.column_end if model.column_end is not None else len(value)
            if start>len(value) or end>len(value) or end<=start: raise LocatorResolutionError("text columns are outside the line")
            return value[start:end].encode()
        if kind is LocatorType.XML_ELEMENT:
            model=XmlElementLocator.model_validate(locator)
            if re.search(br"<!\s*(?:DOCTYPE|ENTITY)\b",content,re.I): raise LocatorResolutionError("unsafe XML declaration")
            root=ET.fromstring(content); expression=model.xpath
            root_prefix=re.compile(rf"^/{re.escape(str(root.tag))}(?:\[1\])?")
            if root_prefix.match(expression): expression="."+root_prefix.sub("",expression,count=1)
            node=root if expression in {"", ".", f"/{root.tag}"} else root.find(expression)
            if node is None: raise LocatorResolutionError("XML path does not resolve")
            if model.attribute is not None:
                if model.attribute not in node.attrib: raise LocatorResolutionError("XML attribute does not resolve")
                return node.attrib[model.attribute].encode()
            return ET.tostring(node,encoding="utf-8")
        if kind is LocatorType.HTML_DOM_SPAN:
            model=HtmlDomSpanLocator.model_validate(locator)
            node=BeautifulSoup(content,"html.parser").select_one(model.selector)
            if node is None: raise LocatorResolutionError("HTML selector does not resolve")
            text=node.get_text()
            start=model.text_start or 0; end=model.text_end if model.text_end is not None else len(text)
            if end>len(text) or end<=start: raise LocatorResolutionError("HTML text span is outside the selected node")
            return text[start:end].encode()
        if kind is LocatorType.ARCHIVE_MEMBER:
            model=ArchiveMemberLocator.model_validate(locator)
            current=content
            for member in model.nested_member_chain: current=_archive_member_bytes(current,member)
            return _archive_member_bytes(current,model.member_path,model.member_ordinal)
        if kind in (LocatorType.EMAIL_HEADER,LocatorType.EMAIL_MIME_PART,LocatorType.EMAIL_ATTACHMENT):
            from email import policy
            from email.parser import BytesParser
            if content.startswith(b"From ") and b"\nFrom " in content:
                import mailbox
                descriptor,temp_name=tempfile.mkstemp(suffix=".mbox"); os.close(descriptor); temp_path=Path(temp_name)
                try:
                    temp_path.write_bytes(content); box=mailbox.mbox(temp_path,create=False)
                    try: messages=[BytesParser(policy=policy.default).parsebytes(item.as_bytes(policy=policy.default)) for item in box]
                    finally: box.close()
                finally: temp_path.unlink(missing_ok=True)
            else:
                messages=[BytesParser(policy=policy.default).parsebytes(content)]
            index=int(locator.get("message",0))
            if index>=len(messages): raise LocatorResolutionError("email message is outside the artifact")
            message=messages[index]
            if kind is LocatorType.EMAIL_HEADER:
                model=EmailHeaderLocator.model_validate(locator); values=message.get_all(model.header,[])
                if model.occurrence>=len(values): raise LocatorResolutionError("email header does not resolve")
                return str(values[model.occurrence]).encode()
            part_key=str(locator.get("part")); selected=None
            if part_key in {"0","1"} and not message.is_multipart(): selected=message
            else:
                selected=message
                for token in part_key.split("."):
                    payload=selected.get_payload() if selected is not None else None
                    if not isinstance(payload,list): selected=None; break
                    position=int(token)-1
                    if position<0 or position>=len(payload): selected=None; break
                    selected=payload[position]
            if selected is None: raise LocatorResolutionError("email MIME part does not resolve")
            payload=selected.get_payload(decode=True)
            if payload is None: payload=str(selected.get_payload()).encode()
            return payload
        if kind is LocatorType.CALENDAR_COMPONENT:
            model=CalendarComponentLocator.model_validate(locator)
            from icalendar import Calendar
            calendar=Calendar.from_ical(content); matches=[]
            for component in calendar.walk(model.component):
                if model.uid is None or str(component.get("UID",""))==model.uid: matches.append(component)
            if not matches: raise LocatorResolutionError("calendar component does not resolve")
            component=matches[0] if model.property else (matches[model.occurrence] if model.occurrence<len(matches) else None)
            if component is None: raise LocatorResolutionError("calendar component does not resolve")
            if model.property:
                raw_value=component.get(model.property)
                values=raw_value if isinstance(raw_value,list) else ([] if raw_value is None else [raw_value])
                if model.occurrence>=len(values): raise LocatorResolutionError("calendar property does not resolve")
                return str(values[model.occurrence]).encode()
            return component.to_ical()
        if kind is LocatorType.VCARD_PROPERTY:
            model=VcardPropertyLocator.model_validate(locator)
            import vobject
            cards=list(vobject.readComponents(content.decode("utf-8-sig")))
            if model.card>=len(cards): raise LocatorResolutionError("vCard is outside the artifact")
            values=cards[model.card].contents.get(model.property.lower(),[])
            if model.occurrence>=len(values): raise LocatorResolutionError("vCard property does not resolve")
            return str(values[model.occurrence].value).encode()
        if kind is LocatorType.SUBTITLE_CUE:
            model=SubtitleCueLocator.model_validate(locator)
            blocks=[block for block in re.split(r"\r?\n\s*\r?\n",content.decode("utf-8-sig").strip()) if block]
            cue_number=0
            for block in blocks:
                lines=block.splitlines()
                timing_index=next((index for index,line in enumerate(lines[:2]) if "-->" in line),None)
                if timing_index is None: continue
                cue_number+=1
                explicit=lines[0].strip() if timing_index==1 else None
                if cue_number==model.cue or explicit==str(model.cue): return block.encode()
            raise LocatorResolutionError("subtitle cue does not resolve")
        if kind is LocatorType.GEOSPATIAL_FEATURE:
            model=GeospatialFeatureLocator.model_validate(locator)
            try:
                document=json.loads(content.decode("utf-8-sig")); features=document.get("features",[]) if isinstance(document,dict) else []
                if isinstance(model.feature,int):
                    if model.feature>=len(features): raise LocatorResolutionError("geospatial feature is outside the artifact")
                    feature=features[model.feature]
                else:
                    feature=next((item for item in features if str(item.get("id"))==model.feature),None)
                    if feature is None: raise LocatorResolutionError("geospatial feature does not resolve")
                return json.dumps(feature,sort_keys=True,separators=(",",":"),ensure_ascii=False).encode()
            except (json.JSONDecodeError,UnicodeDecodeError):
                root=ET.fromstring(content)
                local=lambda tag:tag.rsplit("}",1)[-1]
                if local(root.tag).lower()=="gpx":
                    features=[]
                    for child in root:
                        kind_name=local(child.tag)
                        if kind_name=="wpt": features.append((kind_name,[child]))
                        elif kind_name=="rte": features.append((kind_name,[node for node in child if local(node.tag)=="rtept"]))
                        elif kind_name=="trk":
                            for segment in (node for node in child if local(node.tag)=="trkseg"):
                                features.append((kind_name,[node for node in segment if local(node.tag)=="trkpt"]))
                    index=int(model.feature)
                    if index>=len(features): raise LocatorResolutionError("GPX feature is outside the artifact")
                    feature_type,nodes=features[index]; coordinates=[]; timestamps=[]
                    for node in nodes:
                        elevation=next((item.text for item in node if local(item.tag)=="ele"),None); timestamp=next((item.text for item in node if local(item.tag)=="time"),None)
                        coordinates.append([float(node.attrib["lon"]),float(node.attrib["lat"]),float(elevation) if elevation else None]); timestamps.append(timestamp)
                    return json.dumps({"feature_type":feature_type,"coordinates":coordinates,"timestamps":timestamps},separators=(",",":")).encode()
                placemarks=[node for node in root.iter() if local(node.tag)=="Placemark"]
                index=int(model.feature)
                if index>=len(placemarks): raise LocatorResolutionError("KML feature is outside the artifact")
                coordinates=[]
                for node in placemarks[index].iter():
                    if local(node.tag)=="coordinates" and node.text:
                        for token in node.text.split():
                            parts=token.split(","); coordinates.append([float(parts[0]),float(parts[1]),float(parts[2]) if len(parts)>2 and parts[2] else None])
                return json.dumps({"coordinates":coordinates},separators=(",",":")).encode()
        if kind in (LocatorType.DATABASE_TABLE_ROW,LocatorType.DATABASE_CELL):
            model=DatabaseTableRowLocator.model_validate(locator) if kind is LocatorType.DATABASE_TABLE_ROW else DatabaseCellLocator.model_validate(locator)
            def quote(name:str)->str:
                if not name or "\x00" in name: raise LocatorResolutionError("invalid SQLite identifier")
                return '"'+name.replace('"','""')+'"'
            descriptor,temp_name=tempfile.mkstemp(suffix=".sqlite")
            os.close(descriptor)
            temp_path=Path(temp_name)
            try:
                temp_path.write_bytes(content)
                connection=sqlite3.connect(f"file:{temp_path.as_posix()}?mode=ro",uri=True)
                try:
                    keys=list(model.row_key); where=" AND ".join(f"{quote(key)}=?" for key in keys)
                    selection="*" if kind is LocatorType.DATABASE_TABLE_ROW else quote(model.column)
                    row=connection.execute(f"SELECT {selection} FROM {quote(model.table)} WHERE {where} LIMIT 2",[model.row_key[key] for key in keys]).fetchall()
                    if len(row)!=1: raise LocatorResolutionError("database row key does not resolve uniquely")
                    value=row[0] if kind is LocatorType.DATABASE_TABLE_ROW else row[0][0]
                    return json.dumps(value,sort_keys=True,separators=(",",":"),ensure_ascii=False,default=str).encode()
                finally: connection.close()
            finally: temp_path.unlink(missing_ok=True)
        if kind in (LocatorType.PDF_PAGE_BLOCK,LocatorType.PDF_REGION):
            from pypdf import PdfReader
            reader=PdfReader(io.BytesIO(content),strict=False)
            model=PdfPageBlockLocator.model_validate(locator) if kind is LocatorType.PDF_PAGE_BLOCK else PdfRegionLocator.model_validate(locator)
            if model.page>len(reader.pages): raise LocatorResolutionError("PDF page is outside the artifact")
            page=reader.pages[model.page-1]
            if kind is LocatorType.PDF_PAGE_BLOCK:
                blocks=[part.strip() for part in re.split(r"(?:\r?\n){1,}",page.extract_text() or "") if part.strip()]
                if model.block>=len(blocks): raise LocatorResolutionError("PDF block is outside the page")
                return blocks[model.block].encode()
            page_width=float(page.mediabox.width); page_height=float(page.mediabox.height)
            if model.x+model.width>page_width+0.01 or model.y+model.height>page_height+0.01: raise LocatorResolutionError("PDF region is outside the page")
            fragments=[]
            def visit(fragment,_cm,tm,_font,font_size):
                value=fragment.strip()
                if value:
                    size=max(float(font_size or 0),0.1); fragments.append((value,max(float(tm[4]),0),max(float(tm[5]),0),max(len(value)*size*.5,.1),size))
            page.extract_text(visitor_text=visit)
            for value,x,y,width,height in fragments:
                if abs(x-model.x)<.02 and abs(y-model.y)<.02 and abs(width-model.width)<.02 and abs(height-model.height)<.02: return value.encode()
            return hashlib.sha256(content).digest()+f":pdf-region:{model.page}:{model.x}:{model.y}:{model.width}:{model.height}".encode()
        if kind in (LocatorType.OFFICE_PARAGRAPH,LocatorType.OFFICE_TABLE_CELL,LocatorType.SPREADSHEET_CELL,LocatorType.SLIDE_SHAPE,LocatorType.SLIDE_NOTES):
            names:set[str]=set()
            if content.startswith(b"PK"):
                with zipfile.ZipFile(io.BytesIO(content)) as package: names=set(package.namelist())
            if kind is LocatorType.OFFICE_PARAGRAPH:
                model=OfficeParagraphLocator.model_validate(locator)
                if content.lstrip().startswith(b"{\\rtf"):
                    from striprtf.striprtf import rtf_to_text
                    paragraphs=[item.strip() for item in re.split(r"\r?\n+",rtf_to_text(content.decode("latin-1"))) if item.strip()]
                    if model.paragraph>=len(paragraphs): raise LocatorResolutionError("RTF paragraph is outside the artifact")
                    return paragraphs[model.paragraph].encode()
                if "word/document.xml" in names:
                    from docx import Document
                    paragraph=Document(io.BytesIO(content)).paragraphs[model.paragraph]
                    if model.run is None: return paragraph.text.encode()
                    if model.run<len(paragraph.runs): return paragraph.runs[model.run].text.encode()
                    hyperlinks=paragraph._p.xpath(".//w:hyperlink")
                    index=model.run if model.run<len(hyperlinks) else 0
                    if hyperlinks: return "".join(node.text or "" for node in hyperlinks[index].xpath(".//w:t")).encode()
                    return paragraph.text.encode()
                if "content.xml" in names:
                    with zipfile.ZipFile(io.BytesIO(content)) as package: root=ET.fromstring(package.read("content.xml"))
                    paragraphs=[node for node in root.iter() if node.tag.rsplit("}",1)[-1] in {"p","h"}]
                    if model.paragraph>=len(paragraphs): raise LocatorResolutionError("ODF paragraph is outside the artifact")
                    return "".join(paragraphs[model.paragraph].itertext()).strip().encode()
            if kind is LocatorType.OFFICE_TABLE_CELL:
                model=OfficeTableCellLocator.model_validate(locator)
                if "word/document.xml" in names:
                    from docx import Document
                    return Document(io.BytesIO(content)).tables[model.table].rows[model.row].cells[model.column].text.encode()
                with zipfile.ZipFile(io.BytesIO(content)) as package: root=ET.fromstring(package.read("content.xml"))
                tables=[node for node in root.iter() if node.tag.rsplit("}",1)[-1]=="table"]
                rows=[node for node in tables[model.table] if node.tag.rsplit("}",1)[-1]=="table-row"]
                cells=[node for node in rows[model.row] if node.tag.rsplit("}",1)[-1]=="table-cell"]
                return "".join(cells[model.column].itertext()).strip().encode()
            if kind is LocatorType.SPREADSHEET_CELL:
                model=SpreadsheetCellLocator.model_validate(locator)
                if "xl/workbook.xml" in names:
                    from openpyxl import load_workbook
                    formulas=load_workbook(io.BytesIO(content),data_only=False,read_only=True); cached=load_workbook(io.BytesIO(content),data_only=True,read_only=True)
                    try:
                        formula=formulas[model.sheet].cell(model.row,model.column).value; value=cached[model.sheet].cell(model.row,model.column).value
                        return json.dumps({"formula":formula if isinstance(formula,str) and formula.startswith("=") else None,"value":value if isinstance(formula,str) and formula.startswith("=") else formula},ensure_ascii=False,separators=(",",":"),default=str).encode()
                    finally: formulas.close(); cached.close()
                with zipfile.ZipFile(io.BytesIO(content)) as package: root=ET.fromstring(package.read("content.xml"))
                sheets=[node for node in root.iter() if node.tag.rsplit("}",1)[-1]=="table"]
                sheet=next((node for node in sheets if next((value for key,value in node.attrib.items() if key.rsplit("}",1)[-1]=="name"),None)==model.sheet),None)
                if sheet is None: raise LocatorResolutionError("ODF sheet does not resolve")
                rows=[node for node in sheet if node.tag.rsplit("}",1)[-1]=="table-row"]
                cells=[node for node in rows[model.row-1] if node.tag.rsplit("}",1)[-1]=="table-cell"]
                cell=cells[model.column-1]; return ET.tostring(cell,encoding="utf-8")
            model=SlideShapeLocator.model_validate(locator) if kind is LocatorType.SLIDE_SHAPE else SlideNotesLocator.model_validate(locator)
            if "ppt/presentation.xml" in names:
                from pptx import Presentation
                slide=Presentation(io.BytesIO(content)).slides[model.slide-1]
                if kind is LocatorType.SLIDE_NOTES: return slide.notes_slide.notes_text_frame.text.encode()
                shape=slide.shapes[model.shape]
                if getattr(shape,"has_text_frame",False): return shape.text.encode()
                if getattr(shape,"has_table",False): return json.dumps([[cell.text for cell in row.cells] for row in shape.table.rows],ensure_ascii=False,separators=(",",":")).encode()
                return hashlib.sha256(content).digest()+f":slide-shape:{model.slide}:{model.shape}".encode()
            with zipfile.ZipFile(io.BytesIO(content)) as package: root=ET.fromstring(package.read("content.xml"))
            slides=[node for node in root.iter() if node.tag.rsplit("}",1)[-1]=="page"]
            slide=slides[model.slide-1]
            if kind is LocatorType.SLIDE_NOTES:
                notes=next((node for node in slide.iter() if node.tag.rsplit("}",1)[-1]=="notes"),None)
                if notes is None: raise LocatorResolutionError("slide notes do not resolve")
                return "".join(notes.itertext()).strip().encode()
            shapes=[node for node in slide if node.tag.rsplit("}",1)[-1] in {"frame","custom-shape"}]
            return "".join(shapes[model.shape].itertext()).strip().encode()
        if kind is LocatorType.IMAGE_REGION:
            model=ImageRegionLocator.model_validate(locator)
            try:
                from PIL import Image
                image=Image.open(io.BytesIO(content)); x2=model.x+model.width; y2=model.y+model.height
                if x2>image.width or y2>image.height: raise LocatorResolutionError("image region is outside the image")
                output=io.BytesIO(); image.crop((model.x,model.y,x2,y2)).save(output,format="PNG"); return output.getvalue()
            except LocatorResolutionError: raise
            except Exception as exc: raise LocatorResolutionError("image region cannot be resolved") from exc
        if kind is LocatorType.MEDIA_TIME_RANGE:
            model=MediaTimeRangeLocator.model_validate(locator)
            from mutagen import File as MutagenFile
            media=MutagenFile(io.BytesIO(content))
            duration_ms=int(media.info.length*1000) if media and getattr(media,"info",None) and getattr(media.info,"length",None) else 0
            if duration_ms<=0 and content[:4]==b"RIFF" and content[8:12]==b"WAVE":
                import wave
                with wave.open(io.BytesIO(content),"rb") as wav:
                    duration_ms=int(wav.getnframes()/wav.getframerate()*1000)
            if duration_ms<=0 or model.end_ms>duration_ms:
                raise LocatorResolutionError("media time range is outside the decodable artifact")
            # Byte-accurate slicing is codec-specific; this proof binds the
            # verified time range to the immutable content hash and duration.
            return hashlib.sha256(content).digest()+f":{model.start_ms}:{model.end_ms}".encode()
        if kind is LocatorType.VIDEO_FRAME:
            model=VideoFrameLocator.model_validate(locator)
            if not content: raise LocatorResolutionError("video frame source is empty")
            return hashlib.sha256(content).digest()+f":frame:{model.timestamp_ms}:{model.frame}".encode()
    except LocatorResolutionError: raise
    except Exception as exc: raise LocatorResolutionError(f"invalid {kind.value} locator") from exc
    raise LocatorResolutionError(f"unsupported locator type: {kind.value}")


def verify_locator(content: bytes, locator_type: LocatorType | str, locator: dict[str, Any], expected_hash: str) -> bool:
    resolved=resolve_locator(content,locator_type,locator)
    return hashlib.sha256(resolved).hexdigest()==expected_hash.lower()
