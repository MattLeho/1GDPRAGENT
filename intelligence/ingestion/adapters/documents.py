"""Deterministic extraction for PDF, Office, OpenDocument, and RTF files.

The adapter emits generic document structure and exact canonical locators.  It
does not interpret document meaning and never invokes a provider.  PDF pages
that retain image-only content are merely identified for the ``document.ocr``
task route; OCR output is not substituted for the source bytes here.
"""
from __future__ import annotations

import re
import mimetypes
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from ..models import (
    EmbeddedMember,
    EvidenceLocatorValue,
    ExtractionContext,
    ExtractionResult,
    ExtractionUnit,
    FileTypeTruth,
    ProbeResult,
    QuarantineStatus,
)


_FORMATS = {"pdf", "docx", "xlsx", "pptx", "odt", "ods", "odp", "rtf"}
_FORMAT_BY_SUFFIX = {f".{name}": name for name in _FORMATS}
_ODF_MIMES = {
    "application/vnd.oasis.opendocument.text": "odt",
    "application/vnd.oasis.opendocument.spreadsheet": "ods",
    "application/vnd.oasis.opendocument.presentation": "odp",
}
_OFFICE_CONTENT_TYPES = {
    "word/": "docx",
    "xl/": "xlsx",
    "ppt/": "pptx",
}
_NS = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
    "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
    "presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
    "xlink": "http://www.w3.org/1999/xlink",
    "meta": "urn:oasis:names:tc:opendocument:xmlns:meta:1.0",
    "dc": "http://purl.org/dc/elements/1.1/",
}


def _locator(locator_type: str, **locator: Any) -> EvidenceLocatorValue:
    return EvidenceLocatorValue(locator_type=locator_type, locator=locator)


def _unit(unit_id: str, unit_type: str, ordinal: int, locator_type: str,
          locator: dict[str, Any], *, text: str | None = None, value: Any = None,
          structured_payload: dict[str, Any] | list[Any] | None = None,
          metadata: dict[str, Any] | None = None, parent: str | None = None) -> ExtractionUnit:
    return ExtractionUnit(
        unit_id=unit_id, unit_type=unit_type, ordinal=ordinal, text=text, value=value,
        structured_payload=structured_payload, metadata=metadata or {},
        evidence_locator=_locator(locator_type, **locator), parent_unit_id=parent,
    )


def _normalise_metadata(values: Any) -> dict[str, Any]:
    if values is None:
        return {}
    items = values.items() if hasattr(values, "items") else ()
    return {str(key).lstrip("/"): str(value) for key, value in items if value is not None}


def _package_format(path: str) -> str | None:
    try:
        with zipfile.ZipFile(path) as package:
            names = set(package.namelist())
            if "mimetype" in names:
                mime = package.read("mimetype").decode("ascii", errors="replace").strip()
                if mime in _ODF_MIMES:
                    return _ODF_MIMES[mime]
            for prefix, format_key in _OFFICE_CONTENT_TYPES.items():
                if any(name.startswith(prefix) for name in names):
                    return format_key
    except (OSError, zipfile.BadZipFile):
        return None
    return None


def _media_members(path: str, prefixes: tuple[str, ...]) -> tuple[EmbeddedMember, ...]:
    members: list[EmbeddedMember] = []
    with zipfile.ZipFile(path) as package:
        for name in package.namelist():
            if not name.endswith("/") and name.startswith(prefixes):
                info = package.getinfo(name)
                members.append(EmbeddedMember(
                    member_path=name, ordinal=len(members), declared_size=info.file_size,
                    media_type=mimetypes.guess_type(name)[0],
                    metadata={"package_path": name, "compressed_size": info.compress_size},
                ))
    return tuple(members)


def _xml_text(element: ET.Element) -> str:
    return "".join(element.itertext()).strip()


def _odf_roots(path: str) -> tuple[ET.Element, ET.Element | None]:
    with zipfile.ZipFile(path) as package:
        content = ET.fromstring(package.read("content.xml"))
        metadata = ET.fromstring(package.read("meta.xml")) if "meta.xml" in package.namelist() else None
    return content, metadata


def _odf_metadata(root: ET.Element | None) -> dict[str, Any]:
    if root is None:
        return {}
    result: dict[str, Any] = {}
    for element in root.iter():
        if element is root or not element.text or not element.text.strip():
            continue
        key = element.tag.rsplit("}", 1)[-1]
        value = element.text.strip()
        if key in result:
            current = result[key]
            result[key] = [*current, value] if isinstance(current, list) else [current, value]
        else:
            result[key] = value
    return result


def _safe_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if hasattr(value, "isoformat"):
        return value.isoformat()
    return str(value)


class DocumentsAdapter:
    adapter_id = "documents"
    adapter_version = "1.0.0"
    family = "documents"
    supported_mime_types = frozenset({
        "application/pdf", "application/rtf", "text/rtf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        *_ODF_MIMES,
    })
    supported_extensions = frozenset(_FORMAT_BY_SUFFIX)
    supports_streaming = False
    supports_nested_members = True
    locator_types = frozenset({
        "pdf_page_block", "pdf_region", "office_paragraph", "office_table_cell",
        "spreadsheet_cell", "slide_shape", "slide_notes",
    })
    capability_flags = frozenset({"text", "tables", "embedded_media", "metadata", "pages", "sheets", "slides"})

    def probe(self, path: str, truth: FileTypeTruth) -> ProbeResult:
        try:
            sample = Path(path).read_bytes()[:16]
        except OSError as exc:
            return ProbeResult(accepted=False, confidence=0, reason=f"unreadable: {exc}")
        if sample.startswith(b"%PDF-"):
            return ProbeResult(accepted=True, confidence=1, detected_format="pdf", reason="PDF signature")
        if sample.lstrip().startswith(b"{\\rtf"):
            return ProbeResult(accepted=True, confidence=1, detected_format="rtf", reason="RTF control word")
        package_format = _package_format(path) if sample.startswith(b"PK") else None
        if package_format:
            return ProbeResult(accepted=True, confidence=1, detected_format=package_format, reason="package content signature")
        truth_format = (truth.detected_format or "").lower()
        if truth_format in _FORMATS and truth_format == _FORMAT_BY_SUFFIX.get(Path(path).suffix.lower()):
            return ProbeResult(accepted=False, confidence=0, detected_format=truth_format, reason="declared format lacks required signature")
        return ProbeResult(accepted=False, confidence=0, reason="no supported document signature")

    def extract(self, path: str, context: ExtractionContext) -> ExtractionResult:
        format_key = str(context.configuration.get("detected_format") or self._detect(path)).lower()
        try:
            if format_key == "pdf":
                return self._pdf(path, context)
            if format_key == "docx":
                return self._docx(path, context)
            if format_key == "xlsx":
                return self._xlsx(path, context)
            if format_key == "pptx":
                return self._pptx(path, context)
            if format_key in {"odt", "ods", "odp"}:
                return self._odf(path, context, format_key)
            if format_key == "rtf":
                return self._rtf(path, context)
            raise ValueError(f"unsupported document format: {format_key}")
        except Exception as exc:  # malformed document packages must remain catalogued
            return ExtractionResult(
                artifact_id=context.artifact_id, adapter_id=self.adapter_id,
                adapter_version=self.adapter_version, family=self.family,
                detected_format=format_key or "unknown", warnings=(f"extraction failed: {exc}",),
                quarantine_status=QuarantineStatus.CORRUPT,
            )

    @staticmethod
    def _detect(path: str) -> str:
        sample = Path(path).read_bytes()[:16]
        if sample.startswith(b"%PDF-"):
            return "pdf"
        if sample.lstrip().startswith(b"{\\rtf"):
            return "rtf"
        return _package_format(path) or _FORMAT_BY_SUFFIX.get(Path(path).suffix.lower(), "unknown")

    def _pdf(self, path: str, context: ExtractionContext) -> ExtractionResult:
        from pypdf import PdfReader

        reader = PdfReader(path, strict=False)
        if reader.is_encrypted:
            try:
                unlocked = reader.decrypt("")
            except Exception:
                unlocked = 0
            if not unlocked:
                return ExtractionResult(
                    artifact_id=context.artifact_id, adapter_id=self.adapter_id,
                    adapter_version=self.adapter_version, family=self.family,
                    detected_format="pdf", metadata={"encrypted": True},
                    warnings=("encrypted PDF requires a password",),
                    quarantine_status=QuarantineStatus.PASSWORD_REQUIRED,
                )
        units: list[ExtractionUnit] = []
        members: list[EmbeddedMember] = []
        pages: list[dict[str, Any]] = []
        ocr_pages: list[int] = []
        ordinal = 0
        for page_number, page in enumerate(reader.pages, start=1):
            fragments: list[dict[str, Any]] = []

            def visit_text(fragment: str, _cm: list[float], tm: list[float], _font: Any, font_size: float) -> None:
                content = fragment.strip()
                if not content:
                    return
                size = max(float(font_size or 0), 0.1)
                fragments.append({
                    "text": content, "x": max(float(tm[4]), 0.0), "y": max(float(tm[5]), 0.0),
                    "width": max(len(content) * size * 0.5, 0.1), "height": size,
                })

            text = page.extract_text(visitor_text=visit_text) or ""
            blocks = [part.strip() for part in re.split(r"(?:\r?\n){1,}", text) if part.strip()]
            try:
                images = list(page.images)
            except Exception:
                images = []
            classification = "hybrid" if blocks and images else "scanned_image" if images else "native_text"
            if classification in {"scanned_image", "hybrid"}:
                ocr_pages.append(page_number)
            page_info = {
                "page": page_number, "classification": classification,
                "text_block_count": len(blocks), "embedded_image_count": len(images),
                "width": float(page.mediabox.width), "height": float(page.mediabox.height),
            }
            pages.append(page_info)
            for block_index, block in enumerate(blocks):
                units.append(_unit(
                    f"pdf-page-{page_number}-block-{block_index}", "pdf_text_block", ordinal,
                    "pdf_page_block", {"page": page_number, "block": block_index}, text=block,
                    metadata={"page_classification": classification},
                ))
                ordinal += 1
            for region_index, fragment in enumerate(fragments):
                units.append(_unit(
                    f"pdf-page-{page_number}-region-{region_index}", "pdf_text_region", ordinal,
                    "pdf_region", {"page": page_number, "x": fragment["x"], "y": fragment["y"],
                                   "width": fragment["width"], "height": fragment["height"]},
                    text=fragment["text"], metadata={"coordinate_system": "pdf_user_space", "coordinates_available": True},
                ))
                ordinal += 1
            table_rows = [re.split(r"\t+| {2,}", line.strip()) for line in text.splitlines() if line.strip()]
            if len(table_rows) >= 2 and max((len(row) for row in table_rows), default=0) >= 2:
                units.append(_unit(
                    f"pdf-page-{page_number}-table-0", "pdf_table_candidate", ordinal,
                    "pdf_region", {"page":page_number,"x":0,"y":0,"width":page_info["width"],"height":page_info["height"]},
                    structured_payload={"rows": table_rows}, metadata={"candidate": True},
                ))
                ordinal += 1
            for image_index, image in enumerate(images):
                name = getattr(image, "name", None) or f"page-{page_number}-image-{image_index}"
                data = getattr(image, "data", b"") or b""
                members.append(EmbeddedMember(
                    member_path=f"pdf/page-{page_number}/{name}", ordinal=len(members),
                    declared_size=len(data), media_type=mimetypes.guess_type(name)[0],
                    metadata={"page": page_number, "image": image_index, "name": name},
                ))
                units.append(_unit(
                    f"pdf-page-{page_number}-image-{image_index}", "pdf_embedded_image", ordinal,
                    "pdf_region", {"page":page_number,"x":0,"y":0,"width":page_info["width"],"height":page_info["height"]},
                    structured_payload={"member_path": f"pdf/page-{page_number}/{name}", "size": len(data)},
                    metadata={"coordinates_available": False},
                ))
                ordinal += 1
        metadata = {
            "page_count": len(reader.pages), "pages": pages,
            "document_metadata": _normalise_metadata(reader.metadata),
            "ocr_residue": {"task_route": "document.ocr", "pages": ocr_pages},
            "raw_source_preserved": True,
        }
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format="pdf",
            metadata=metadata, units=tuple(units), embedded_members=tuple(members),
        )

    def _docx(self, path: str, context: ExtractionContext) -> ExtractionResult:
        from docx import Document

        document = Document(path)
        units: list[ExtractionUnit] = []
        ordinal = 0
        for paragraph_index, paragraph in enumerate(document.paragraphs):
            text = paragraph.text
            style = paragraph.style.name if paragraph.style is not None else None
            if text:
                unit_type = "heading" if style and style.lower().startswith("heading") else "paragraph"
                units.append(_unit(
                    f"docx-paragraph-{paragraph_index}", unit_type, ordinal, "office_paragraph",
                    {"paragraph": paragraph_index}, text=text, metadata={"style": style},
                ))
                ordinal += 1
            hyperlink_index = 0
            for hyperlink in paragraph._p.xpath(".//w:hyperlink"):
                target_id = hyperlink.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                target = document.part.rels[target_id].target_ref if target_id in document.part.rels else None
                link_text = "".join(node.text or "" for node in hyperlink.xpath(".//w:t"))
                units.append(_unit(
                    f"docx-paragraph-{paragraph_index}-link-{hyperlink_index}", "hyperlink", ordinal,
                    "office_paragraph", {"paragraph": paragraph_index, "run": hyperlink_index},
                    text=link_text or target or "", structured_payload={"target": target},
                    parent=f"docx-paragraph-{paragraph_index}" if text else None,
                ))
                ordinal += 1
                hyperlink_index += 1
            for comment_index, marker in enumerate(paragraph._p.xpath(".//w:commentRangeStart")):
                raw_id = marker.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id")
                try:
                    comment = document.comments.get(int(raw_id))
                except (TypeError, ValueError):
                    comment = None
                if comment is not None:
                    units.append(_unit(
                        f"docx-paragraph-{paragraph_index}-comment-{raw_id}", "comment", ordinal,
                        "office_paragraph", {"paragraph": paragraph_index, "run": comment_index},
                        text=comment.text, metadata={"comment_id": int(raw_id), "author": comment.author,
                                                     "initials": comment.initials, "timestamp": _safe_value(comment.timestamp)},
                        parent=f"docx-paragraph-{paragraph_index}" if text else None,
                    ))
                    ordinal += 1
        for table_index, table in enumerate(document.tables):
            for row_index, row in enumerate(table.rows):
                for column_index, cell in enumerate(row.cells):
                    units.append(_unit(
                        f"docx-table-{table_index}-{row_index}-{column_index}", "table_cell", ordinal,
                        "office_table_cell", {"table": table_index, "row": row_index, "column": column_index},
                        text=cell.text, metadata={"table": table_index},
                    ))
                    ordinal += 1
        properties = document.core_properties
        metadata = {name: _safe_value(getattr(properties, name, None)) for name in (
            "title", "subject", "author", "keywords", "comments", "created", "modified", "last_modified_by", "revision"
        ) if getattr(properties, name, None) is not None}
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format="docx",
            metadata={"document_metadata": metadata, "paragraph_count": len(document.paragraphs), "table_count": len(document.tables)},
            units=tuple(units), embedded_members=_media_members(path, ("word/media/",)),
        )

    def _xlsx(self, path: str, context: ExtractionContext) -> ExtractionResult:
        from openpyxl import load_workbook

        formulas = load_workbook(path, data_only=False, read_only=False)
        cached = load_workbook(path, data_only=True, read_only=False)
        units: list[ExtractionUnit] = []
        sheets: list[dict[str, Any]] = []
        ordinal = 0
        try:
            for sheet in formulas.worksheets:
                cached_sheet = cached[sheet.title]
                sheet_info = {
                    "name": sheet.title, "state": sheet.sheet_state,
                    "used_range": sheet.calculate_dimension(),
                    "merged_ranges": [str(item) for item in sheet.merged_cells.ranges],
                    "table_regions": [str(table.ref) for table in sheet.tables.values()],
                }
                sheets.append(sheet_info)
                for row in sheet.iter_rows():
                    for cell in row:
                        formula = cell.value if cell.data_type == "f" else None
                        raw_value = None if formula is not None else cell.value
                        cached_value = cached_sheet[cell.coordinate].value if formula is not None else cell.value
                        if cell.value is None and cell.comment is None:
                            continue
                        payload = {
                            "sheet": sheet.title, "row": cell.row, "column": cell.column,
                            "address": cell.coordinate, "formula": formula,
                            "value": _safe_value(raw_value), "cached_value": _safe_value(cached_value),
                            "display_value": _safe_value(cached_value),
                            "value_source": "cached_formula_result" if formula is not None and cached_value is not None else "formula_without_cached_result" if formula is not None else "cell_value",
                        }
                        if cell.comment is not None:
                            payload["comment"] = {"text": cell.comment.text, "author": cell.comment.author}
                        units.append(_unit(
                            f"xlsx-{sheet.title}-{cell.coordinate}", "spreadsheet_cell", ordinal,
                            "spreadsheet_cell", {"sheet": sheet.title, "row": cell.row, "column": cell.column, "address": cell.coordinate},
                            value=_safe_value(cached_value if formula is not None else cell.value), structured_payload=payload,
                            metadata={"sheet_state": sheet.sheet_state},
                        ))
                        ordinal += 1
        finally:
            formulas.close()
            cached.close()
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format="xlsx",
            metadata={"sheet_names": [sheet["name"] for sheet in sheets], "sheets": sheets},
            units=tuple(units), embedded_members=_media_members(path, ("xl/media/",)),
        )

    def _pptx(self, path: str, context: ExtractionContext) -> ExtractionResult:
        from pptx import Presentation

        presentation = Presentation(path)
        units: list[ExtractionUnit] = []
        slides: list[dict[str, Any]] = []
        ordinal = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = slide.shapes.title.text if slide.shapes.title is not None else None
            slides.append({"slide": slide_number, "title": title, "shape_count": len(slide.shapes)})
            for shape_index, shape in enumerate(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.text:
                    units.append(_unit(
                        f"pptx-slide-{slide_number}-shape-{shape_index}", "slide_text_shape", ordinal,
                        "slide_shape", {"slide": slide_number, "shape": shape_index}, text=shape.text,
                        metadata={"name": shape.name, "shape_type": str(shape.shape_type), "is_title": shape is slide.shapes.title},
                    ))
                    ordinal += 1
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    units.append(_unit(
                        f"pptx-slide-{slide_number}-table-{shape_index}", "slide_table", ordinal,
                        "slide_shape", {"slide": slide_number, "shape": shape_index},
                        structured_payload={"rows": rows}, metadata={"name": shape.name},
                    ))
                    ordinal += 1
                if getattr(shape, "has_chart", False):
                    units.append(_unit(
                        f"pptx-slide-{slide_number}-chart-{shape_index}", "slide_chart", ordinal,
                        "slide_shape", {"slide": slide_number, "shape": shape_index},
                        structured_payload={"chart_type": str(shape.chart.chart_type), "series_count": len(shape.chart.series)},
                        metadata={"name": shape.name},
                    ))
                    ordinal += 1
            try:
                notes_text = presentation.slides[slide_number - 1].notes_slide.notes_text_frame.text
            except (AttributeError, KeyError):
                notes_text = ""
            if notes_text.strip():
                units.append(_unit(
                    f"pptx-slide-{slide_number}-notes", "slide_notes", ordinal,
                    "slide_notes", {"slide": slide_number, "note": 0}, text=notes_text,
                ))
                ordinal += 1
        properties = presentation.core_properties
        metadata = {name: _safe_value(getattr(properties, name, None)) for name in (
            "title", "subject", "author", "keywords", "comments", "created", "modified", "last_modified_by", "revision"
        ) if getattr(properties, name, None) is not None}
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format="pptx",
            metadata={"slide_count": len(presentation.slides), "slides": slides, "document_metadata": metadata},
            units=tuple(units), embedded_members=_media_members(path, ("ppt/media/",)),
        )

    def _odf(self, path: str, context: ExtractionContext, format_key: str) -> ExtractionResult:
        root, metadata_root = _odf_roots(path)
        if format_key == "odt":
            units, details = self._odt_units(root)
        elif format_key == "ods":
            units, details = self._ods_units(root)
        else:
            units, details = self._odp_units(root)
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format=format_key,
            metadata={"document_metadata": _odf_metadata(metadata_root), **details}, units=tuple(units),
            embedded_members=_media_members(path, ("Pictures/",)),
        )

    def _odt_units(self, root: ET.Element) -> tuple[list[ExtractionUnit], dict[str, Any]]:
        units: list[ExtractionUnit] = []
        ordinal = 0
        paragraphs = list(root.findall(".//text:p", _NS)) + list(root.findall(".//text:h", _NS))
        for paragraph_index, paragraph in enumerate(paragraphs):
            text = _xml_text(paragraph)
            if text:
                heading = paragraph.tag == f"{{{_NS['text']}}}h"
                units.append(_unit(
                    f"odt-paragraph-{paragraph_index}", "heading" if heading else "paragraph", ordinal,
                    "office_paragraph", {"paragraph": paragraph_index}, text=text,
                    metadata={"outline_level": paragraph.get(f"{{{_NS['text']}}}outline-level") if heading else None},
                ))
                ordinal += 1
            for run_index, link in enumerate(paragraph.findall(".//text:a", _NS)):
                target = link.get(f"{{{_NS['xlink']}}}href")
                units.append(_unit(
                    f"odt-paragraph-{paragraph_index}-link-{run_index}", "hyperlink", ordinal,
                    "office_paragraph", {"paragraph": paragraph_index, "run": run_index},
                    text=_xml_text(link) or target or "", structured_payload={"target": target},
                ))
                ordinal += 1
            for comment_index, annotation in enumerate(paragraph.findall(".//office:annotation", _NS)):
                comment_text = _xml_text(annotation)
                if comment_text:
                    units.append(_unit(
                        f"odt-paragraph-{paragraph_index}-comment-{comment_index}", "comment", ordinal,
                        "office_paragraph", {"paragraph": paragraph_index, "run": comment_index},
                        text=comment_text,
                    ))
                    ordinal += 1
        tables = root.findall(".//table:table", _NS)
        for table_index, table in enumerate(tables):
            for row_index, row in enumerate(table.findall("table:table-row", _NS)):
                for column_index, cell in enumerate(row.findall("table:table-cell", _NS)):
                    units.append(_unit(
                        f"odt-table-{table_index}-{row_index}-{column_index}", "table_cell", ordinal,
                        "office_table_cell", {"table": table_index, "row": row_index, "column": column_index},
                        text=_xml_text(cell),
                    ))
                    ordinal += 1
        return units, {"paragraph_count": len(paragraphs), "table_count": len(tables)}

    def _ods_units(self, root: ET.Element) -> tuple[list[ExtractionUnit], dict[str, Any]]:
        units: list[ExtractionUnit] = []
        sheets: list[dict[str, Any]] = []
        ordinal = 0
        for sheet in root.findall(".//table:table", _NS):
            name = sheet.get(f"{{{_NS['table']}}}name") or f"Sheet{len(sheets) + 1}"
            visibility = sheet.get(f"{{{_NS['table']}}}display", "true")
            row_number = 1
            maximum_column = 0
            for row in sheet.findall("table:table-row", _NS):
                row_repeat = min(int(row.get(f"{{{_NS['table']}}}number-rows-repeated", "1")), 10000)
                column_number = 1
                for cell in row.findall("table:table-cell", _NS):
                    repeat = min(int(cell.get(f"{{{_NS['table']}}}number-columns-repeated", "1")), 10000)
                    text = _xml_text(cell)
                    formula = cell.get(f"{{{_NS['table']}}}formula")
                    value = next((cell.get(f"{{{_NS['office']}}}{kind}") for kind in ("value", "date-value", "time-value", "boolean-value", "string-value") if cell.get(f"{{{_NS['office']}}}{kind}") is not None), None)
                    if text or formula or value is not None:
                        for offset in range(repeat):
                            column = column_number + offset
                            address = self._cell_address(row_number, column)
                            payload = {"sheet": name, "row": row_number, "column": column, "address": address,
                                       "formula": formula, "value": value, "cached_value": text or value,
                                       "display_value": text, "value_source": "cached_formula_result" if formula else "cell_value"}
                            units.append(_unit(
                                f"ods-{name}-{address}", "spreadsheet_cell", ordinal, "spreadsheet_cell",
                                {"sheet": name, "row": row_number, "column": column, "address": address},
                                value=text or value, structured_payload=payload, metadata={"sheet_state": "visible" if visibility != "false" else "hidden"},
                            ))
                            ordinal += 1
                    column_number += repeat
                maximum_column = max(maximum_column, column_number - 1)
                row_number += row_repeat
            sheets.append({"name": name, "state": "visible" if visibility != "false" else "hidden", "used_rows": row_number - 1, "used_columns": maximum_column})
        return units, {"sheet_names": [sheet["name"] for sheet in sheets], "sheets": sheets}

    def _odp_units(self, root: ET.Element) -> tuple[list[ExtractionUnit], dict[str, Any]]:
        units: list[ExtractionUnit] = []
        ordinal = 0
        pages = root.findall(".//draw:page", _NS)
        for slide_number, page in enumerate(pages, start=1):
            shape_index = 0
            notes = page.find("presentation:notes", _NS)
            for shape in list(page):
                if shape is notes:
                    continue
                text = _xml_text(shape)
                if text:
                    units.append(_unit(
                        f"odp-slide-{slide_number}-shape-{shape_index}", "slide_text_shape", ordinal,
                        "slide_shape", {"slide": slide_number, "shape": shape_index}, text=text,
                        metadata={"name": shape.get(f"{{{_NS['draw']}}}name")},
                    ))
                    ordinal += 1
                shape_index += 1
            for table_index, table in enumerate(page.findall(".//table:table", _NS)):
                rows = [[_xml_text(cell) for cell in row.findall("table:table-cell", _NS)]
                        for row in table.findall("table:table-row", _NS)]
                units.append(_unit(
                    f"odp-slide-{slide_number}-table-{table_index}", "slide_table", ordinal,
                    "slide_shape", {"slide": slide_number, "shape": shape_index + table_index},
                    structured_payload={"rows": rows},
                ))
                ordinal += 1
            if notes is not None and _xml_text(notes):
                units.append(_unit(
                    f"odp-slide-{slide_number}-notes", "slide_notes", ordinal,
                    "slide_notes", {"slide": slide_number, "note": 0}, text=_xml_text(notes),
                ))
                ordinal += 1
        return units, {"slide_count": len(pages)}

    def _rtf(self, path: str, context: ExtractionContext) -> ExtractionResult:
        from striprtf.striprtf import rtf_to_text

        raw = Path(path).read_text(encoding="latin-1")
        text = rtf_to_text(raw)
        paragraphs = [paragraph.strip() for paragraph in re.split(r"\r?\n+", text) if paragraph.strip()]
        units = [
            _unit(f"rtf-paragraph-{index}", "paragraph", index, "office_paragraph", {"paragraph": index}, text=paragraph)
            for index, paragraph in enumerate(paragraphs)
        ]
        return ExtractionResult(
            artifact_id=context.artifact_id, adapter_id=self.adapter_id,
            adapter_version=self.adapter_version, family=self.family, detected_format="rtf",
            metadata={"paragraph_count": len(paragraphs)}, units=tuple(units),
        )

    @staticmethod
    def _cell_address(row: int, column: int) -> str:
        letters = ""
        current = column
        while current:
            current, remainder = divmod(current - 1, 26)
            letters = chr(65 + remainder) + letters
        return f"{letters}{row}"
