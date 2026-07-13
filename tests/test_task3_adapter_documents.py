from __future__ import annotations

import io
import base64
import zipfile
from pathlib import Path
from uuid import uuid4

import pytest
from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches as PptxInches
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from ingestion.adapters.documents import DocumentsAdapter
from ingestion.models import ExtractionContext, FileTypeTruth, FileTypeTruthValue


_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def _context(path: Path, detected_format: str | None = None) -> ExtractionContext:
    return ExtractionContext(
        artifact_id=uuid4(), analysis_run_id=uuid4(), export_snapshot_id=uuid4(),
        source_path=str(path), configuration={"detected_format": detected_format} if detected_format else {},
    )


def _truth(detected_format: str | None = None) -> FileTypeTruth:
    return FileTypeTruth(
        status=FileTypeTruthValue.UNKNOWN, detected_format=detected_format,
        evidence=(), reason="test",
    )


def _pdf(path: Path, page_kinds: tuple[str, ...]) -> None:
    output = canvas.Canvas(str(path), pagesize=(200, 200))
    image = ImageReader(io.BytesIO(_PNG))
    for index, kind in enumerate(page_kinds, start=1):
        if kind in {"native_text", "hybrid"}:
            output.drawString(20, 170, f"Native page {index}")
        if kind in {"scanned_image", "hybrid"}:
            output.drawImage(image, 20, 20, width=160, height=130, mask="auto")
        output.showPage()
    output.save()


def test_pdf_classifies_native_scanned_and_hybrid_and_routes_only_residue(tmp_path: Path):
    path = tmp_path / "mixed.pdf"
    _pdf(path, ("native_text", "scanned_image", "hybrid"))

    adapter = DocumentsAdapter()
    assert adapter.probe(str(path), _truth()).detected_format == "pdf"
    result = adapter.extract(str(path), _context(path))

    assert [page["classification"] for page in result.metadata["pages"]] == [
        "native_text", "scanned_image", "hybrid",
    ]
    assert result.metadata["ocr_residue"] == {"task_route": "document.ocr", "pages": [2, 3]}
    assert result.metadata["raw_source_preserved"] is True
    assert {unit.evidence_locator.locator_type for unit in result.units} == {"pdf_page_block", "pdf_region"}
    assert all(unit.metadata["coordinates_available"] for unit in result.units if unit.unit_type == "pdf_text_region")
    assert {member.metadata["page"] for member in result.embedded_members} == {2, 3}


def test_pdf_native_table_candidate_and_malformed_pdf_are_bounded(tmp_path: Path):
    native = tmp_path / "native.pdf"
    output = canvas.Canvas(str(native), pagesize=(200, 200))
    output.drawString(20, 170, "Name  Value")
    output.drawString(20, 150, "Alpha  10")
    output.save()
    result = DocumentsAdapter().extract(str(native), _context(native))
    assert result.metadata["page_count"] == 1
    assert any(unit.unit_type == "pdf_text_block" for unit in result.units)

    malformed = tmp_path / "broken.pdf"
    malformed.write_bytes(b"%PDF-1.7\nnot a PDF")
    broken = DocumentsAdapter().extract(str(malformed), _context(malformed))
    assert broken.quarantine_status.value == "corrupt"
    assert broken.warnings


def test_docx_preserves_paragraph_heading_table_and_embedded_media(tmp_path: Path):
    image = tmp_path / "pixel.png"
    image.write_bytes(_PNG)
    path = tmp_path / "document.docx"
    document = Document()
    document.core_properties.title = "Fixture document"
    document.add_heading("Evidence heading", level=1)
    paragraph = document.add_paragraph("A generic paragraph")
    document.add_comment(paragraph.runs, text="Review note", author="Reviewer", initials="RR")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    document.add_picture(str(image), width=Inches(0.2))
    document.save(path)

    result = DocumentsAdapter().extract(str(path), _context(path))

    assert result.metadata["document_metadata"]["title"] == "Fixture document"
    assert any(unit.unit_type == "heading" and unit.text == "Evidence heading" for unit in result.units)
    comment = next(unit for unit in result.units if unit.unit_type == "comment")
    assert comment.text == "Review note"
    assert comment.metadata["author"] == "Reviewer"
    cells = [unit for unit in result.units if unit.unit_type == "table_cell"]
    assert [cell.evidence_locator.locator for cell in cells] == [
        {"table": 0, "row": 0, "column": 0}, {"table": 0, "row": 0, "column": 1},
    ]
    assert result.embedded_members[0].member_path.startswith("word/media/")


def test_xlsx_preserves_formula_cached_distinction_merged_cells_hidden_sheet_comments_and_media(tmp_path: Path):
    image = tmp_path / "pixel.png"
    image.write_bytes(_PNG)
    path = tmp_path / "workbook.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Visible"
    sheet["A1"] = 2
    sheet["B1"] = 3
    sheet["C1"] = "=SUM(A1:B1)"
    sheet["A2"] = "Merged"
    sheet.merge_cells("A2:B2")
    sheet["A1"].comment = Comment("Cell note", "Tester")
    sheet.add_image(__import__("openpyxl").drawing.image.Image(str(image)), "E1")
    hidden = workbook.create_sheet("Hidden")
    hidden.sheet_state = "hidden"
    hidden["A1"] = "secret"
    workbook.save(path)

    result = DocumentsAdapter().extract(str(path), _context(path))

    visible = next(item for item in result.metadata["sheets"] if item["name"] == "Visible")
    assert visible["merged_ranges"] == ["A2:B2"]
    assert next(item for item in result.metadata["sheets"] if item["name"] == "Hidden")["state"] == "hidden"
    formula = next(unit for unit in result.units if unit.evidence_locator.locator["address"] == "C1")
    assert formula.structured_payload["formula"] == "=SUM(A1:B1)"
    assert formula.structured_payload["cached_value"] is None
    assert formula.structured_payload["value_source"] == "formula_without_cached_result"
    comment = next(unit for unit in result.units if unit.evidence_locator.locator["address"] == "A1")
    assert comment.structured_payload["comment"] == {"text": "Cell note", "author": "Tester"}
    assert result.embedded_members[0].member_path.startswith("xl/media/")


def test_pptx_preserves_slide_shapes_table_speaker_notes_and_media(tmp_path: Path):
    image = tmp_path / "pixel.png"
    image.write_bytes(_PNG)
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide title"
    table = slide.shapes.add_table(2, 2, PptxInches(1), PptxInches(1), PptxInches(4), PptxInches(1)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    slide.shapes.add_picture(str(image), PptxInches(1), PptxInches(3), width=PptxInches(0.2))
    slide.notes_slide.notes_text_frame.text = "Speaker evidence note"
    presentation.save(path)

    result = DocumentsAdapter().extract(str(path), _context(path))

    assert result.metadata["slide_count"] == 1
    assert any(unit.unit_type == "slide_table" for unit in result.units)
    note = next(unit for unit in result.units if unit.unit_type == "slide_notes")
    assert "Speaker evidence note" in note.text
    assert note.evidence_locator.locator == {"slide": 1, "note": 0}
    assert all(unit.evidence_locator.locator["slide"] == 1 for unit in result.units)
    assert result.embedded_members[0].member_path.startswith("ppt/media/")


def _odf(path: Path, mime: str, content: str, media: bool = False) -> None:
    with zipfile.ZipFile(path, "w") as package:
        package.writestr("mimetype", mime, compress_type=zipfile.ZIP_STORED)
        package.writestr("content.xml", content)
        package.writestr(
            "meta.xml",
            "<office:document-meta xmlns:office='urn:oasis:names:tc:opendocument:xmlns:office:1.0' "
            "xmlns:dc='http://purl.org/dc/elements/1.1/'><office:meta><dc:title>ODF fixture</dc:title></office:meta></office:document-meta>",
        )
        if media:
            package.writestr("Pictures/pixel.png", _PNG)


@pytest.mark.parametrize(
    "suffix,mime,content,locator_type",
    [
        ("odt", "application/vnd.oasis.opendocument.text",
         "<office:document-content xmlns:office='urn:oasis:names:tc:opendocument:xmlns:office:1.0' xmlns:text='urn:oasis:names:tc:opendocument:xmlns:text:1.0'><office:body><office:text><text:h text:outline-level='1'>Heading</text:h><text:p>Paragraph</text:p></office:text></office:body></office:document-content>", "office_paragraph"),
        ("ods", "application/vnd.oasis.opendocument.spreadsheet",
         "<office:document-content xmlns:office='urn:oasis:names:tc:opendocument:xmlns:office:1.0' xmlns:table='urn:oasis:names:tc:opendocument:xmlns:table:1.0' xmlns:text='urn:oasis:names:tc:opendocument:xmlns:text:1.0'><office:body><office:spreadsheet><table:table table:name='Data'><table:table-row><table:table-cell office:value-type='float' office:value='5'><text:p>5</text:p></table:table-cell><table:table-cell table:formula='of:=SUM([.A1])' office:value-type='float' office:value='5'><text:p>5</text:p></table:table-cell></table:table-row></table:table></office:spreadsheet></office:body></office:document-content>", "spreadsheet_cell"),
        ("odp", "application/vnd.oasis.opendocument.presentation",
         "<office:document-content xmlns:office='urn:oasis:names:tc:opendocument:xmlns:office:1.0' xmlns:draw='urn:oasis:names:tc:opendocument:xmlns:drawing:1.0' xmlns:presentation='urn:oasis:names:tc:opendocument:xmlns:presentation:1.0' xmlns:text='urn:oasis:names:tc:opendocument:xmlns:text:1.0'><office:body><office:presentation><draw:page><draw:frame><text:p>Slide text</text:p></draw:frame><presentation:notes><text:p>Slide note</text:p></presentation:notes></draw:page></office:presentation></office:body></office:document-content>", "slide_shape"),
    ],
)
def test_open_document_formats_use_canonical_structural_locators(tmp_path: Path, suffix: str, mime: str, content: str, locator_type: str):
    path = tmp_path / f"fixture.{suffix}"
    _odf(path, mime, content, media=True)

    adapter = DocumentsAdapter()
    assert adapter.probe(str(path), _truth()).detected_format == suffix
    result = adapter.extract(str(path), _context(path))

    assert result.metadata["document_metadata"]["title"] == "ODF fixture"
    assert any(unit.evidence_locator.locator_type == locator_type for unit in result.units)
    assert result.embedded_members[0].member_path == "Pictures/pixel.png"
    if suffix == "ods":
        formula = next(unit for unit in result.units if unit.evidence_locator.locator["address"] == "B1")
        assert formula.structured_payload["formula"] == "of:=SUM([.A1])"
        assert formula.structured_payload["cached_value"] == "5"
    if suffix == "odp":
        assert any(unit.evidence_locator.locator_type == "slide_notes" for unit in result.units)


def test_rtf_paragraphs_and_signature_probe(tmp_path: Path):
    path = tmp_path / "document.rtf"
    path.write_bytes(b"{\\rtf1\\ansi First paragraph\\par Second paragraph\\par}")
    adapter = DocumentsAdapter()

    assert adapter.probe(str(path), _truth()).detected_format == "rtf"
    result = adapter.extract(str(path), _context(path))
    assert [unit.text for unit in result.units] == ["First paragraph", "Second paragraph"]
    assert [unit.evidence_locator.locator for unit in result.units] == [{"paragraph": 0}, {"paragraph": 1}]


def test_extension_alone_never_accepts_a_document(tmp_path: Path):
    path = tmp_path / "fake.docx"
    path.write_bytes(b"plain text with the wrong extension")
    probe = DocumentsAdapter().probe(str(path), _truth("docx"))
    assert probe.accepted is False
    assert "signature" in probe.reason
